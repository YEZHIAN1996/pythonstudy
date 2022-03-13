import pptx
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

p = pptx.Presentation()  # 生成ppt对象
layout = p.slide_layouts[1]  # 选择布局，有0-8九种，一般常使用1和7
slide = p.slides.add_slide(layout)  # 创建一页ppt

title = slide.placeholders[0]
content = slide.placeholders[1]

title.text = '题目'
# content.text = 'python写ppt学习'
paragraph1 = content.text_frame.add_paragraph()
paragraph1.text = '欢迎学习ppt制作'
paragraph1.bold = True
paragraph1.font.italic = True
paragraph1.font.underline = True
paragraph1.font.size = Pt(16)
paragraph1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

paragraph2 = content.text_frame.add_paragraph()
paragraph2.text = '欢迎学习python'
paragraph2.font.size = Pt(32)
paragraph2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

layout1 = p.slide_layouts[6]  # title
slide2 = p.slides.add_slide(layout1)
top = left = width = height = Inches(5)
box = slide2.shapes.add_textbox(left, top, width, height)
para = box.text_frame.add_paragraph()
para.text = 'this is a para test'
para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
para.font.size = Pt(32)
para.font.color.rgb = RGBColor(255, 255, 0)
para.font.name = '微软雅黑'

layout2 = p.slide_layouts[1]
slide3 = p.slides.add_slide(layout2)
rows = 10
cols = 2
left = top = Inches(2)
width = Inches(6.0)
height = Inches(1.0)

table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
for index, _ in enumerate(range(rows)):
    for sub_index in range(cols):
        table.cell(index, sub_index).text = '%s:%s' % (index,sub_index)


layout3 = p.slide_layouts[6]
slide4 = p.slides.add_slide(layout3)
image = slide4.shapes.add_picture('logo2020.png', left, top, width, height)



p.save('a.pptx')